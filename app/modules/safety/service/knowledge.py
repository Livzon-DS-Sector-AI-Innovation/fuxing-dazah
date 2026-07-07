"""Safety knowledge article business workflows."""

import logging
import os
import uuid
from datetime import date, datetime

from fastapi import UploadFile
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.storage import delete_object
from app.core.storage import is_enabled as minio_enabled
from app.modules.safety.models import (
    SafetyKnowledgeArticle,
)
from app.modules.safety.repository import SafetyRepository
from app.modules.safety.schemas import (
    AgentUsageStatsResponse,
    BatchGenerateCardsResponse,
    DuplicateArticleItem,
    DuplicateCheckResponse,
    GenerateCardResponse,
    GeneratePptResponse,
    GenerateSummaryResponse,
    KnowledgeCardContent,
    KnowledgeSyncResponse,
    PptGenerationRecord,
    PptHistoryResponse,
    SafetyKnowledgeArticleCreate,
    SafetyKnowledgeArticleParseResponse,
    SafetyKnowledgeArticleUpdate,
    VersionChainItem,
)

logger = logging.getLogger(__name__)

# 类别 → 编号前缀映射
CATEGORY_PREFIX_MAP: dict[str, str] = {
    "laws_regulations": "LAW",
    "standards": "STD",
    "management_systems": "MGT",
    "accident_cases": "CASE",
    "emergency_plans": "ERP",
    "sds": "SDS",
    "training_materials": "TRN",
    "other": "GEN",
}


class KnowledgeService:
    """安全知识库业务服务"""

    def __init__(self, session: AsyncSession):
        self.session = session
        self.repo = SafetyRepository(session)

    # ── 文件清理 ──────────────────────────────────────────

    @staticmethod
    def _cleanup_file(file_path: str | None) -> None:
        """Delete a single file from MinIO or local disk."""
        if not file_path:
            return
        try:
            if minio_enabled():
                try:
                    delete_object("safety", file_path)
                except Exception:
                    pass
            else:
                abs_path = os.path.abspath(file_path)
                if os.path.exists(abs_path):
                    os.remove(abs_path)
        except OSError:
            pass

    # ── 编号生成 ──────────────────────────────────────────

    def _generate_article_no(self, category: str) -> str:
        """生成文档编号：{PREFIX}-{YYYYMMDD}-{3位序号}"""
        prefix = CATEGORY_PREFIX_MAP.get(category, "GEN")
        today = date.today().strftime("%Y%m%d")
        # Note: max_seq query requires await; handled in create_article
        return f"{prefix}-{today}"

    # ── CRUD ──────────────────────────────────────────────

    async def get_articles(
        self,
        skip: int = 0,
        limit: int = 20,
        category: str | None = None,
        status: str | None = None,
        keyword: str | None = None,
    ) -> tuple[list[SafetyKnowledgeArticle], int]:
        """获取知识库文章列表"""
        return await self.repo.get_knowledge_articles(
            skip, limit, category, status, keyword
        )

    async def get_article(self, article_id: uuid.UUID) -> SafetyKnowledgeArticle | None:
        """获取文章详情（浏览计数+1）"""
        article = await self.repo.get_knowledge_article_by_id(article_id)
        if article:
            await self.repo.update_knowledge_article(
                article_id, {"view_count": article.view_count + 1}
            )
        return article

    async def create_article(
        self, data: SafetyKnowledgeArticleCreate
    ) -> SafetyKnowledgeArticle:
        """创建知识库文章（含自动编号）"""
        article_data = data.model_dump()

        # Auto-generate article_no if not provided
        if not article_data.get("article_no"):
            prefix = self._generate_article_no(data.category.value)
            max_seq = await self.repo.get_max_article_seq_for_date(
                prefix, date.today().strftime("%Y%m%d")
            )
            seq = max_seq + 1
            article_data["article_no"] = f"{prefix}-{seq:03d}"

        return await self.repo.create_knowledge_article(article_data)

    async def update_article(
        self, article_id: uuid.UUID, data: SafetyKnowledgeArticleUpdate
    ) -> SafetyKnowledgeArticle | None:
        """更新知识库文章"""
        update_data = {k: v for k, v in data.model_dump().items() if v is not None}
        if not update_data:
            return None
        return await self.repo.update_knowledge_article(article_id, update_data)

    async def delete_article(self, article_id: uuid.UUID) -> bool:
        """删除知识库文章（软删除 + 清理附件）"""
        article = await self.repo.get_knowledge_article_by_id(article_id)
        result = await self.repo.delete_knowledge_article(article_id)
        if result and article:
            self._cleanup_file(article.attachment_path)
        return result

    async def publish_article(self, article_id: uuid.UUID) -> SafetyKnowledgeArticle | None:
        """发布文章（草稿→已发布）"""
        article = await self.repo.get_knowledge_article_by_id(article_id)
        if not article or article.status != "draft":
            return None
        return await self.repo.update_knowledge_article(
            article_id, {"status": "published"}
        )

    async def archive_article(self, article_id: uuid.UUID) -> SafetyKnowledgeArticle | None:
        """归档文章（已发布→已归档）"""
        article = await self.repo.get_knowledge_article_by_id(article_id)
        if not article or article.status != "published":
            return None
        return await self.repo.update_knowledge_article(
            article_id, {"status": "archived"}
        )

    # ── AI 智能解析 ───────────────────────────────────────

    async def parse_document_metadata(
        self, file: UploadFile
    ) -> SafetyKnowledgeArticleParseResponse:
        """AI 解析上传的文档，提取元数据。

        Steps:
        1. Save uploaded file to temp location
        2. Extract text via DocumentParser
        3. Call AI to extract structured metadata
        4. Return parse result (not persisted to DB)
        """
        from app.platform.integrations.ai.document_parser import DocumentParser

        # Save temp file for parsing
        content_bytes = await file.read()
        file_ext = os.path.splitext(file.filename or ".txt")[1].lower()
        temp_dir = os.path.join("uploads", "safety", "knowledge", "_temp")
        os.makedirs(temp_dir, exist_ok=True)
        temp_path = os.path.join(
            temp_dir,
            f"parse_{uuid.uuid4().hex[:8]}_{datetime.now().timestamp()}{file_ext}",
        )

        try:
            with open(temp_path, "wb") as f:
                f.write(content_bytes)

            # Extract full text
            full_text = DocumentParser.extract_text(temp_path, max_chars=30000)

            if not full_text or len(full_text.strip()) < 20:
                raise ValueError("无法从文件中提取有效文本内容，请确认文件格式正确")

            # AI extraction of metadata
            result = await self._ai_extract_metadata(full_text)

            return result

        finally:
            # Clean up temp file
            try:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
            except OSError:
                pass

    async def _ai_extract_metadata(self, full_text: str) -> SafetyKnowledgeArticleParseResponse:
        """Call AI to extract structured metadata from document text."""
        from app.modules.safety.service.config import create_ai_service

        # Truncate for AI input (keep prompt manageable)
        text_for_ai = full_text[:8000]
        if len(full_text) > 8000:
            text_for_ai += "\n\n（文档过长，已截断前8000字符用于分析）"

        prompt = f"""你是一个安全法规文档分析专家。请从以下文档原文中提取关键元数据。

## 文档原文
{text_for_ai}

## 提取要求
请以 JSON 格式返回以下字段，所有判断必须基于原文内容：

1. **title**: 文档的完整标题（如"GB 3836.1-2021 爆炸性环境 第1部分：设备 通用要求"）
2. **category**: 知识分类，必须是以下之一：
   - laws_regulations: 法律法规（如安全生产法、危化品管理条例、部门规章）
   - standards: 标准规范（如 GB、GB/T、HG、SH 等标准）
   - management_systems: 安全管理制度（企业内部制度、操作规程汇编）
   - accident_cases: 事故案例
   - emergency_plans: 应急预案
   - sds: 化学品安全技术说明书
   - training_materials: 培训教材
   - other: 其他
3. **summary**: 文档核心内容摘要（100-200字），涵盖主要章节和关键规定
4. **tags**: 关键标签（逗号分隔，5-10个），如"防爆,电气安全,危险区域划分,GB3836"
5. **source**: 文档来源/发布机构（如"国家标准化管理委员会"、"应急管理部"）
6. **author**: 作者/起草单位（如有）
7. **publish_date**: 发布日期（YYYY-MM-DD 格式，如能从原文推断）

请直接返回 JSON，不要包含任何其他文字：
{{"title": "...", "category": "...", "summary": "...", "tags": "...", "source": "...", "author": "...", "publish_date": "..."}}"""

        ai_service = await create_ai_service("text")
        try:
            parsed = await ai_service.chat_parsed(
                messages=[{"role": "user", "content": prompt}],
                expected_keys=[
                    "title", "category", "summary", "tags",
                    "source", "author", "publish_date",
                ],
            )

            # Validate category
            valid_categories = [
                "laws_regulations", "standards", "management_systems",
                "accident_cases", "emergency_plans", "sds",
                "training_materials", "other",
            ]
            category = parsed.get("category", "other")
            if category not in valid_categories:
                category = "other"

            return SafetyKnowledgeArticleParseResponse(
                title=parsed.get("title", "未命名文档"),
                summary=parsed.get("summary", ""),
                tags=parsed.get("tags", ""),
                source=parsed.get("source", ""),
                author=parsed.get("author", ""),
                publish_date=parsed.get("publish_date"),
                content_preview=full_text[:500],
                full_content=full_text,
            )
        finally:
            await ai_service.close()

    async def batch_parse_documents(
        self, files: list[UploadFile]
    ) -> list[SafetyKnowledgeArticleParseResponse]:
        """批量解析多个文档文件"""
        results: list[SafetyKnowledgeArticleParseResponse] = []
        for file in files:
            try:
                result = await self.parse_document_metadata(file)
                # Attach original filename
                results.append(result)
            except Exception as e:
                logger.warning("批量解析文件 %s 失败: %s", file.filename, e)
        return results

    # ── AI 知识卡片生成 ───────────────────────────────────

    async def generate_knowledge_card(
        self, article_id: uuid.UUID
    ) -> GenerateCardResponse | None:
        """AI 从文档全文生成 6 维度知识卡片。

        流程：
        1. 读取 article.content（全文）
        2. 构建结构化 prompt
        3. 调用 AI 生成 6 维度内容
        4. 写入 article.knowledge_card + card_version++
        5. 返回卡片内容
        """
        from app.modules.safety.service.config import create_ai_service

        article = await self.repo.get_knowledge_article_by_id(article_id)
        if not article:
            return None

        content = article.content or ""
        if not content or len(content.strip()) < 50:
            return GenerateCardResponse(
                knowledge_card=KnowledgeCardContent(),
                card_version=article.card_version or 1,
                message="文档正文内容不足，无法生成知识卡片。请先确保文档已解析全文。",
            )

        # Build generation prompt
        text_for_ai = content[:12000]
        if len(content) > 12000:
            text_for_ai += "\n\n（文档过长，已截断前12000字符用于分析）"

        prompt = f"""你是一个安全法规与标准文档分析专家。请从以下文档原文中提取结构化知识，生成知识卡片。

## 文档原文
{text_for_ai}

## 知识卡片 6 维度说明

请按以下 6 个维度提取核心内容。每个维度只提取原文中直接相关的内容，**不要编造**。如果文档未涉及某维度，该字段留空字符串。

1. **hazard_type_definitions** — 危险源类型定义：文档中定义的危险源类型、分类方式，特别是人/物/环/管四类危险源的具体描述
2. **hazard_category_criteria** — 隐患分类标准：如何判定隐患属于哪一类别（如电气、机械、火灾爆炸等13类）的准则
3. **hazard_level_criteria** — 隐患分级标准：重大隐患/一般隐患/低风险的判定依据和定量/定性标准
4. **key_defect_examples** — 典型缺陷示例：文档中列出的具体安全缺陷、不合规情形范例
5. **rectification_requirements** — 整改措施要求：对不同等级隐患的整改时限、整改标准、防护要求
6. **legal_basis_clauses** — 法律依据条文：可直接引用的法规条款原文（条款号 + 原文）

## 输出格式
请严格返回 JSON，每个字段值为提取出的原文摘要（200-500字），无相关内容则留空字符串：

{{"hazard_type_definitions": "...", "hazard_category_criteria": "...", "hazard_level_criteria": "...", "key_defect_examples": "...", "rectification_requirements": "...", "legal_basis_clauses": "..."}}"""

        ai_service = await create_ai_service("text")
        try:
            parsed = await ai_service.chat_parsed(
                messages=[{"role": "user", "content": prompt}],
                expected_keys=[
                    "hazard_type_definitions",
                    "hazard_category_criteria",
                    "hazard_level_criteria",
                    "key_defect_examples",
                    "rectification_requirements",
                    "legal_basis_clauses",
                ],
            )

            card_content = KnowledgeCardContent(
                hazard_type_definitions=parsed.get("hazard_type_definitions", "") or None,
                hazard_category_criteria=parsed.get("hazard_category_criteria", "") or None,
                hazard_level_criteria=parsed.get("hazard_level_criteria", "") or None,
                key_defect_examples=parsed.get("key_defect_examples", "") or None,
                rectification_requirements=parsed.get("rectification_requirements", "") or None,
                legal_basis_clauses=parsed.get("legal_basis_clauses", "") or None,
            )

            new_version = (article.card_version or 0) + 1

            # Save to article
            await self.repo.update_knowledge_article(
                article_id,
                {
                    "knowledge_card": card_content.model_dump(),
                    "card_version": new_version,
                },
            )

            return GenerateCardResponse(
                knowledge_card=card_content,
                card_version=new_version,
                message="知识卡片生成成功",
            )

        finally:
            await ai_service.close()

    async def get_agent_usage_stats(
        self, article_id: uuid.UUID
    ) -> AgentUsageStatsResponse | None:
        """获取该文档的知识卡片被 Agent 注入的使用统计。

        当前为占位实现（后续可接入注入日志表）。
        """
        article = await self.repo.get_knowledge_article_by_id(article_id)
        if not article:
            return None

        return AgentUsageStatsResponse(
            article_id=article.id,
            article_title=article.title or "",
            total_injections_30d=0,
            by_agent={
                "hazard_identification": 0,
                "hazard_source_identification": 0,
                "rectification_review": 0,
            },
            last_injected_at=None,
        )

    async def batch_generate_cards(
        self, article_ids: list[uuid.UUID]
    ) -> BatchGenerateCardsResponse:
        """批量生成知识卡片（顺序执行，单条失败不影响其他）。"""
        results: list[dict] = []
        success_count = 0
        failed_count = 0

        for aid in article_ids:
            try:
                result = await self.generate_knowledge_card(aid)
                if result and result.knowledge_card.hazard_type_definitions:
                    success_count += 1
                    results.append({"id": str(aid), "success": True, "message": result.message})
                elif result:
                    failed_count += 1
                    results.append({"id": str(aid), "success": False, "message": result.message})
                else:
                    failed_count += 1
                    results.append({"id": str(aid), "success": False, "message": "文档不存在"})
            except Exception as e:
                failed_count += 1
                logger.warning("批量生成卡片失败 article_id=%s: %s", aid, e)
                results.append({"id": str(aid), "success": False, "message": str(e)})

        return BatchGenerateCardsResponse(
            success_count=success_count,
            failed_count=failed_count,
            results=results,
        )

    # ── AI PPT 生成 ───────────────────────────────────────

    async def generate_ppt(
        self,
        article_id: uuid.UUID,
        template: str = "training",
        style: str = "professional",
    ) -> GeneratePptResponse | None:
        """AI 从文档内容生成培训 PPT。

        流程：
        1. 读取 article.content + knowledge_card
        2. 构建 PPT 大纲 prompt
        3. 调用 AI 生成每页标题 + 内容
        4. 用 python-pptx 生成 .pptx 文件
        5. 保存到 uploads 目录，返回下载 URL
        """
        article = await self.repo.get_knowledge_article_by_id(article_id)
        if not article:
            return None

        content = article.content or ""
        if not content or len(content.strip()) < 100:
            return GeneratePptResponse(
                download_url="",
                file_name="",
                page_count=0,
                message="文档正文内容不足（至少需要100字），无法生成 PPT。",
            )

        # Build card context if available
        card_context = ""
        if article.knowledge_card:
            card = article.knowledge_card
            card_parts = []
            field_labels = {
                "hazard_type_definitions": "危险源类型定义",
                "hazard_category_criteria": "隐患分类标准",
                "hazard_level_criteria": "隐患分级标准",
                "key_defect_examples": "典型缺陷示例",
                "rectification_requirements": "整改措施要求",
                "legal_basis_clauses": "法律依据条文",
            }
            for key, label in field_labels.items():
                val = card.get(key)
                if val:
                    card_parts.append(f"### {label}\n{val}")
            if card_parts:
                card_context = "\n\n".join(card_parts)

        # Build AI prompt
        text_for_ai = content[:8000]
        if len(content) > 8000:
            text_for_ai += "\n\n（文档过长，已截断前8000字符）"

        template_descriptions = {
            "training": "员工安全培训课件 — 适合用于部门安全培训、新员工入职安全教育",
            "briefing": "安全简报 — 适合管理层汇报、安全形势分析",
            "audit": "审核检查清单 — 适合安全检查对照、合规性审计",
        }
        style_descriptions = {
            "professional": "专业蓝白配色，正式严谨",
            "modern": "现代简约深色背景，适合投影演示",
            "minimal": "极简白底黑字，适合打印分发",
        }

        prompt = f"""你是一个安全培训 PPT 制作专家。请根据以下文档内容，制作一份结构化 PPT 大纲。

## 文档内容
{text_for_ai}

## 知识卡片参考
{card_context if card_context else "（无知识卡片）"}

## PPT 要求
- 模板类型: {template_descriptions.get(template, template)}
- 配色风格: {style_descriptions.get(style, style)}
- 页数: 10-20 页

## 输出格式
请严格返回 JSON，结构如下：
{{
  "title": "PPT 总标题",
  "subtitle": "副标题（可选）",
  "slides": [
    {{"title": "第1页标题", "content": ["要点1", "要点2", "要点3"], "notes": "演讲备注"}},
    ...
  ]
}}

每页 slide 的 content 数组包含 3-6 个要点，使用简洁的短句。
第一页为封面（title + subtitle），最后一页为总结/致谢。
请直接返回 JSON，不要包含任何其他文字。"""

        from app.modules.safety.service.config import create_ai_service

        ai_service = await create_ai_service("text")
        try:
            parsed = await ai_service.chat_parsed(
                messages=[{"role": "user", "content": prompt}],
                expected_keys=["title", "slides"],
            )

            ppt_title = parsed.get("title", article.title or "安全培训课件")
            ppt_subtitle = parsed.get("subtitle", "")
            slides = parsed.get("slides", [])

            if not slides:
                return GeneratePptResponse(
                    download_url="",
                    file_name="",
                    page_count=0,
                    message="AI 未能生成有效的 PPT 大纲，请检查文档内容质量。",
                )

            # Generate PPTX file
            file_name, file_path = self._build_pptx_file(
                ppt_title, ppt_subtitle, slides, template, style
            )

            # Build relative path for file proxy
            file_path_rel = os.path.join("uploads", "safety", "knowledge", "ppt", file_name)

            return GeneratePptResponse(
                download_url=file_path_rel,
                file_name=file_name,
                page_count=len(slides),
                message=f"PPT 生成成功，共 {len(slides)} 页",
            )

        finally:
            await ai_service.close()

    @staticmethod
    def _build_pptx_file(
        title: str,
        subtitle: str,
        slides: list[dict],
        template: str,
        style: str,
    ) -> tuple[str, str]:
        """用 python-pptx 生成 .pptx 文件。

        Returns:
            (file_name, abs_path) 元组
        """
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Style presets
        style_presets = {
            "professional": {
                "bg": RGBColor(0xFF, 0xFF, 0xFF),
                "title_color": RGBColor(0x1A, 0x3C, 0x6E),
                "body_color": RGBColor(0x33, 0x33, 0x33),
                "accent": RGBColor(0x2B, 0x5A, 0xA7),
            },
            "modern": {
                "bg": RGBColor(0x1E, 0x1E, 0x2E),
                "title_color": RGBColor(0xFF, 0xFF, 0xFF),
                "body_color": RGBColor(0xD0, 0xD0, 0xD0),
                "accent": RGBColor(0x00, 0xB4, 0xD8),
            },
            "minimal": {
                "bg": RGBColor(0xFF, 0xFF, 0xFF),
                "title_color": RGBColor(0x22, 0x22, 0x22),
                "body_color": RGBColor(0x55, 0x55, 0x55),
                "accent": RGBColor(0x66, 0x66, 0x66),
            },
        }
        preset = style_presets.get(style, style_presets["professional"])

        for i, slide_data in enumerate(slides):
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = preset["bg"]

            slide_title = slide_data.get("title", "")
            slide_content = slide_data.get("content", [])
            slide_notes = slide_data.get("notes", "")

            if i == 0:
                # Cover slide
                # Title
                title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = preset["title_color"]
                p.alignment = PP_ALIGN.CENTER

                # Subtitle
                if subtitle:
                    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(1))
                    tf2 = sub_box.text_frame
                    p2 = tf2.paragraphs[0]
                    p2.text = subtitle
                    p2.font.size = Pt(24)
                    p2.font.color.rgb = preset["body_color"]
                    p2.alignment = PP_ALIGN.CENTER

                # Accent line
                line = slide.shapes.add_shape(
                    1,  # MSO_SHAPE.RECTANGLE
                    Inches(5), Inches(3.5), Inches(3.333), Pt(3),
                )
                line.fill.solid()
                line.fill.fore_color.rgb = preset["accent"]
                line.line.fill.background()
            else:
                # Title bar
                title_bar = slide.shapes.add_shape(
                    1, Inches(0), Inches(0), Inches(13.333), Inches(1.2),
                )
                title_bar.fill.solid()
                title_bar.fill.fore_color.rgb = preset["accent"]
                title_bar.line.fill.background()

                title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.733), Inches(0.9))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                # Content
                content_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.333), Inches(5.2))
                tf2 = content_box.text_frame
                tf2.word_wrap = True

                for j, bullet in enumerate(slide_content):
                    if j == 0:
                        p2 = tf2.paragraphs[0]
                    else:
                        p2 = tf2.add_paragraph()
                    p2.text = f"• {bullet}"
                    p2.font.size = Pt(20)
                    p2.font.color.rgb = preset["body_color"]
                    p2.space_after = Pt(12)

                # Page number
                page_box = slide.shapes.add_textbox(Inches(11.5), Inches(7), Inches(1.5), Inches(0.4))
                tf3 = page_box.text_frame
                p3 = tf3.paragraphs[0]
                p3.text = f"{i + 1}/{len(slides)}"
                p3.font.size = Pt(10)
                p3.font.color.rgb = preset["body_color"]
                p3.alignment = PP_ALIGN.RIGHT

            # Speaker notes
            if slide_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_notes

        # Save file
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c for c in title if c.isalnum() or c in "._- ")[:50].strip()
        file_name = f"ppt_{safe_title}_{ts}.pptx"
        ppt_dir = os.path.join("uploads", "safety", "knowledge", "ppt")
        os.makedirs(ppt_dir, exist_ok=True)
        file_path = os.path.join(ppt_dir, file_name)
        prs.save(file_path)

        return file_name, file_path

    async def get_ppt_history(self, article_id: uuid.UUID) -> PptHistoryResponse:
        """获取文档的 PPT 生成历史记录。（从文件系统扫描）"""
        ppt_dir = os.path.join("uploads", "safety", "knowledge", "ppt")
        if not os.path.isdir(ppt_dir):
            return PptHistoryResponse(records=[])

        records: list[PptGenerationRecord] = []
        article = await self.repo.get_knowledge_article_by_id(article_id)
        if not article or not article.title:
            return PptHistoryResponse(records=[])

        safe_title = "".join(c for c in article.title if c.isalnum() or c in "._- ")[:50].strip()
        prefix = f"ppt_{safe_title}_"

        try:
            for entry in sorted(
                os.scandir(ppt_dir),
                key=lambda e: e.stat().st_mtime,
                reverse=True,
            ):
                if entry.is_file() and entry.name.startswith(prefix) and entry.name.endswith(".pptx"):
                    stat = entry.stat()
                    rel_path = os.path.join("uploads", "safety", "knowledge", "ppt", entry.name)
                    records.append(
                        PptGenerationRecord(
                            id=entry.name.replace(".pptx", ""),
                            file_name=entry.name,
                            template="training",
                            style="professional",
                            page_count=0,  # Would need to read file to get count
                            download_url=rel_path,
                            created_at=datetime.fromtimestamp(stat.st_mtime).isoformat(),
                        )
                    )
        except OSError:
            pass

        return PptHistoryResponse(records=records)

    # ── AI 摘要生成 ───────────────────────────────────────

    async def generate_summary(
        self, article_id: uuid.UUID
    ) -> GenerateSummaryResponse | None:
        """AI 从文档全文生成结构化摘要。

        流程：
        1. 读取 article.content
        2. 构建摘要提取 prompt
        3. 调用 AI 生成摘要
        4. 回写到 article.summary
        5. 返回摘要
        """
        article = await self.repo.get_knowledge_article_by_id(article_id)
        if not article:
            return None

        content = article.content or ""
        if not content or len(content.strip()) < 50:
            return GenerateSummaryResponse(
                summary="",
                message="文档正文内容不足，无法生成摘要。",
            )

        text_for_ai = content[:10000]
        if len(content) > 10000:
            text_for_ai += "\n\n（文档过长，已截断前10000字符）"

        prompt = f"""你是一个安全法规文档分析专家。请为以下文档生成结构化摘要。

## 文档内容
{text_for_ai}

## 摘要要求
请生成一份 150-300 字的结构化摘要，涵盖以下要素：
1. **文档类型**：是什么类型的文档（法规/标准/制度/案例等）
2. **核心主题**：文档的核心内容是什么
3. **关键要点**：3-5 个最重要的规定或发现
4. **适用范围**：适用于哪些场景/部门/人员

请直接返回摘要文本（纯文本，不要 Markdown 标题），语言简洁专业。"""

        from app.modules.safety.service.config import create_ai_service

        ai_service = await create_ai_service("text")
        try:
            parsed = await ai_service.chat_parsed(
                messages=[{"role": "user", "content": prompt}],
                expected_keys=["summary"],
            )

            summary = parsed.get("summary", "")
            if not summary:
                # Fallback: treat entire response as summary
                summary = parsed.get("_raw", "")[:500] if isinstance(parsed, dict) else str(parsed)[:500]

            if not summary:
                return GenerateSummaryResponse(
                    summary="",
                    message="AI 未能生成有效摘要，请检查文档内容质量。",
                )

            # Save to article
            await self.repo.update_knowledge_article(
                article_id,
                {"summary": summary},
            )

            return GenerateSummaryResponse(
                summary=summary,
                message="摘要生成成功，已保存到文档。",
            )

        finally:
            await ai_service.close()

    # ── 重复检测 ──────────────────────────────────────────

    async def check_duplicate(
        self, title: str, content: str | None = None
    ) -> DuplicateCheckResponse:
        """检测是否有相似文档。

        先用 ILIKE 做标题模糊匹配快速筛查，有 content 时调用 AI 做语义相似度判断。
        """
        # Step 1: Quick ILIKE filter
        similar = await self.repo.check_similar_articles(title)

        if not similar:
            return DuplicateCheckResponse(has_duplicates=False, duplicates=[])

        # Step 2: If we have content, use AI for semantic similarity
        if content and len(similar) > 0:
            duplicates = await self._ai_check_similarity(title, content, similar)
        else:
            # Without content, just return ILIKE results as potential duplicates
            duplicates = [
                DuplicateArticleItem(
                    id=a.id,
                    article_no=a.article_no,
                    title=a.title,
                    category=a.category,
                    similarity_reason="标题相似",
                )
                for a in similar
            ]

        return DuplicateCheckResponse(
            has_duplicates=len(duplicates) > 0,
            duplicates=duplicates,
        )

    async def _ai_check_similarity(
        self,
        title: str,
        content: str,
        candidates: list[SafetyKnowledgeArticle],
    ) -> list[DuplicateArticleItem]:
        """用 AI 判断新文档与候选文档的语义相似度。"""
        from app.modules.safety.service.config import create_ai_service

        candidates_text = "\n".join(
            f"[{i}] 标题: {c.title} | 摘要: {(c.summary or '')[:200]}"
            for i, c in enumerate(candidates)
        )

        prompt = f"""你是一个文档去重专家。请判断以下新文档是否与知识库中已有文档高度重复（实质内容相同）。

## 新文档
标题: {title}
内容: {(content or '')[:1000]}

## 候选已有文档
{candidates_text}

## 判断要求
- 如果两篇文档论述的是同一个标准/法规/主题，且主要内容重叠超过70%，视为重复
- 同一法规的不同版本（如 GB 3836.1-2010 vs GB 3836.1-2021）视为非重复（不同版本）
- 不同标准即使涉及同一主题也不算重复
- 返回 JSON: {{"duplicates": [0, 2]}} 其中数字是候选文档的索引
- 如果没有重复，返回 {{"duplicates": []}}"""

        ai_service = await create_ai_service("text")
        try:
            parsed = await ai_service.chat_parsed(
                messages=[{"role": "user", "content": prompt}],
                expected_keys=["duplicates"],
            )
            dup_indices = parsed.get("duplicates", [])
            if isinstance(dup_indices, list):
                return [
                    DuplicateArticleItem(
                        id=candidates[i].id,
                        article_no=candidates[i].article_no,
                        title=candidates[i].title,
                        category=candidates[i].category,
                        similarity_reason="AI 判断内容高度相似",
                    )
                    for i in dup_indices
                    if isinstance(i, int) and 0 <= i < len(candidates)
                ]
            return []
        except Exception as e:
            logger.warning("AI相似度检测失败: %s，回退为标题匹配", e)
            return [
                DuplicateArticleItem(
                    id=a.id,
                    article_no=a.article_no,
                    title=a.title,
                    category=a.category,
                    similarity_reason="标题相似（AI检测降级）",
                )
                for a in candidates[:3]
            ]
        finally:
            await ai_service.close()

    # ── 版本管理 ──────────────────────────────────────────

    async def get_version_chain(
        self, article_id: uuid.UUID
    ) -> list[VersionChainItem]:
        """获取文档的完整版本链。"""
        article = await self.repo.get_knowledge_article_by_id(article_id)
        if not article:
            return []
        return await self._build_version_chain(article)

    async def create_new_version(
        self, article_id: uuid.UUID
    ) -> tuple[SafetyKnowledgeArticle | None, list[VersionChainItem]]:
        """基于现有文档创建新版本。

        Returns:
            (new_article, version_chain) — 新版本草案 + 完整版本链
        """
        original = await self.repo.get_knowledge_article_by_id(article_id)
        if not original:
            return None, []

        new_version_num = (original.version or 1) + 1

        # Create new draft based on original
        from app.modules.safety.schemas import SafetyKnowledgeArticleCreate

        new_data = SafetyKnowledgeArticleCreate(
            article_no=original.article_no,
            title=original.title,
            summary=original.summary,
            content=original.content,
            tags=original.tags,
            source=original.source,
            author=original.author,
            publish_date=original.publish_date,
            notes=original.notes,
            category=original.category,
        )

        new_article = await self.repo.create_knowledge_article(
            {
                **new_data.model_dump(),
                "version": new_version_num,
                "status": "draft",
            }
        )

        # Mark original as superseded by new version
        await self.repo.update_knowledge_article(
            article_id, {"superseded_by_id": new_article.id}
        )

        # Build version chain
        version_chain = await self._build_version_chain(new_article)

        return new_article, version_chain

    async def _build_version_chain(
        self, article: SafetyKnowledgeArticle
    ) -> list[VersionChainItem]:
        """构建版本链（从旧到新排列）。"""
        from sqlalchemy import select as sa_select

        visited: set[uuid.UUID] = set()

        # Start from the current article, find all connected versions
        # via superseded_by_id chain
        all_versions: list[SafetyKnowledgeArticle] = [article]
        visited.add(article.id)

        # Find predecessors: documents whose superseded_by_id is in our list
        # Loop to build complete chain
        changed = True
        while changed:
            changed = False
            superseded_ids = [a.id for a in all_versions]
            stmt = sa_select(SafetyKnowledgeArticle).where(
                SafetyKnowledgeArticle.superseded_by_id.in_(superseded_ids),
                not SafetyKnowledgeArticle.is_deleted,
            )
            result = await self.session.execute(stmt)
            predecessors = list(result.scalars().all())
            for p in predecessors:
                if p.id not in visited:
                    visited.add(p.id)
                    all_versions.append(p)
                    changed = True

        # Sort by version ascending (oldest first), then by created_at
        all_versions.sort(
            key=lambda a: (a.version or 1, a.created_at or datetime.min)
        )

        # Find the latest version (highest version number)
        latest_version = max(a.version or 1 for a in all_versions)

        chain: list[VersionChainItem] = []
        for a in all_versions:
            chain.append(
                VersionChainItem(
                    id=a.id,
                    article_no=a.article_no,
                    title=a.title,
                    version=a.version or 1,
                    status=a.status,
                    is_current=(a.version or 1) == latest_version,
                    created_at=a.created_at or datetime.min,
                )
            )

        return chain

    # ── 语义搜索 ──────────────────────────────────────────

    async def semantic_search(
        self,
        query: str,
        skip: int = 0,
        limit: int = 20,
    ) -> tuple[list[dict], int]:
        """自然语言语义搜索知识库。

        流程：
        1. AI 解析查询意图 → 提取关键词
        2. 用关键词做 ILIKE 搜索
        3. 返回排序结果
        """
        from app.modules.safety.service.config import create_ai_service

        # Step 1: AI extracts search keywords
        prompt = f"""你是一个安全知识库搜索专家。用户想搜索知识库中的文档，请解析其查询意图。

用户查询: "{query}"

请提取 3-5 个核心搜索关键词（中文），用于模糊匹配文档的标题、摘要、正文和标签。
关键词应当是具体的术语或短语，不要泛化。
返回 JSON: {{"keywords": ["关键词1", "关键词2", ...]}}"""

        ai_service = await create_ai_service("text")
        try:
            parsed = await ai_service.chat_parsed(
                messages=[{"role": "user", "content": prompt}],
                expected_keys=["keywords"],
            )
            keywords: list[str] = parsed.get("keywords", [query])
        except Exception as e:
            logger.warning("语义搜索关键词提取失败: %s，使用原始查询", e)
            keywords = [query]
        finally:
            await ai_service.close()

        # Step 2: Search with extracted keywords
        items, total = await self.repo.search_by_keywords(
            keywords, skip=skip, limit=limit
        )

        # Step 3: Format results
        result_items = []
        for item in items:
            result_items.append({
                "id": item.id,
                "article_no": item.article_no,
                "title": item.title,
                "category": item.category,
                "summary": item.summary,
                "tags": item.tags,
                "status": item.status,
                "match_reason": f"匹配关键词: {', '.join(keywords[:3])}",
            })

        return result_items, total

    # ═══════════════════════════════════════════════════════════════
    # Bitable 全量同步
    # ═══════════════════════════════════════════════════════════════

    async def sync_from_bitable(self) -> KnowledgeSyncResponse:
        """从 Bitable 全量同步知识库（安全管理制度 + 法规标准 + 设备说明书）。

        以 Bitable 为数据源，三阶段同步：
        ① CREATE — Bitable 有、平台无
        ② UPDATE — 双方都有，字段有差异
        ③ DELETE — 平台有、Bitable 无（软删除）
        """
        from sqlalchemy import func, select, update

        from app.modules.safety.feishu.bitable_client import SafetyBitableClient

        # ── 数据表定义 ──
        tables: list[dict[str, object]] = [
            {
                "table_id": "tbl36YLqyB1S5qos",
                "name": "安全管理制度",
                "source": "制度库",
                "field_map": {
                    "制度名称": "title",
                    "制度编号": "article_no",
                    "所属分类": "category",
                    "制度状态": "status",
                    "备注": "notes",
                },
                "category_map": {
                    "目标职责": "management_systems",
                    "制度化管理": "management_systems",
                    "教育培训": "training_materials",
                    "现场管理": "management_systems",
                    "安全风险管控及隐患排查治理": "management_systems",
                    "安全风险管控及隐患排查": "management_systems",
                    "应急管理": "emergency_plans",
                    "事故管理": "accident_cases",
                    "持续改进": "management_systems",
                },
                "status_map": {
                    "现行有效": "published",
                    "拟修订": "draft",
                    "修订中": "draft",
                    "已废止": "archived",
                },
                "category_field": "所属分类",
                "attachment_field": "制度原件",
                "version_field": "版本",
                "article_no_prefix": {
                    "management_systems": "PROC",
                    "training_materials": "TRN",
                    "emergency_plans": "ERP",
                    "accident_cases": "CASE",
                    "laws_regulations": "LAW",
                    "standards": "STD",
                    "other": "GEN",
                },
            },
            {
                "table_id": "tbl85HKWCTfyf6rw",
                "name": "法规标准",
                "source": "法规标准库",
                "field_map": {
                    "法律法规及标准名称": "title",
                    "法规类别": "category",
                    "法规状态": "status",
                    "颁布机关": "author",
                    "核心要点总结": "content",
                    "备注": "notes",
                },
                "category_map": {
                    "安全类": "laws_regulations",
                    "建筑防火与消防": "laws_regulations",
                    "特种设备": "laws_regulations",
                    "特殊作业": "laws_regulations",
                    "危险作业": "laws_regulations",
                    "职业健康": "laws_regulations",
                    "化学品管理": "laws_regulations",
                    "其他相关法规": "laws_regulations",
                },
                "status_map": {
                    "现行有效": "published",
                    "现行有效(新)": "published",
                    "即将实施": "draft",
                    "征求意见中": "draft",
                },
                "category_field": "法规类别",
                "attachment_field": "法规原件",
                "version_field": None,
                "article_no_prefix": {
                    "laws_regulations": "LAW",
                    "standards": "STD",
                    "management_systems": "PROC",
                    "training_materials": "TRN",
                    "emergency_plans": "ERP",
                    "accident_cases": "CASE",
                    "other": "GEN",
                },
            },
        ]

        default_prefix: dict[str, str] = {
            "laws_regulations": "LAW",
            "standards": "STD",
            "management_systems": "PROC",
            "accident_cases": "CASE",
            "emergency_plans": "ERP",
            "sds": "SDS",
            "training_materials": "TRN",
            "other": "GEN",
        }

        # ── 工具函数 ──

        def _extract_rich_text(value: object) -> str:
            if isinstance(value, list):
                parts: list[str] = []
                for item in value:
                    if isinstance(item, dict):
                        parts.append(str(item.get("text", "")))
                    elif item:
                        parts.append(str(item))
                return "".join(parts)
            if isinstance(value, str):
                return value
            return str(value) if value else ""

        def _extract_select_value(value: object) -> str:
            if isinstance(value, str):
                return value
            if isinstance(value, list) and len(value) > 0:
                if isinstance(value[0], dict):
                    return str(value[0].get("text", "") or value[0].get("name", ""))
                return str(value[0])
            if isinstance(value, dict):
                return str(value.get("text", "") or value.get("name", "") or "")
            return ""

        def _extract_attachment_list(value: object) -> list[dict[str, object]]:
            if isinstance(value, list):
                return [
                    {"file_token": a.get("file_token", ""), "name": a.get("name", "")}
                    for a in value if isinstance(a, dict)
                ]
            return []

        def _ms_timestamp_to_date(ms_value: object) -> date | None:
            if not ms_value:
                return None
            try:
                ts = int(str(ms_value)) / 1000
                return datetime.fromtimestamp(ts).date()
            except (ValueError, TypeError, OSError):
                return None

        def map_fields(
            bitable_fields: dict[str, object], table_def: dict[str, object]
        ) -> dict[str, object]:
            result: dict[str, object] = {}
            field_map = table_def["field_map"]
            category_map = table_def["category_map"]
            status_map = table_def["status_map"]
            category_field = str(table_def["category_field"])
            attachment_field = str(table_def["attachment_field"])
            version_field = table_def["version_field"]

            for cn_name, en_name in field_map.items():
                raw = bitable_fields.get(cn_name)
                if raw is None or raw == "" or raw == []:
                    continue

                if cn_name == category_field:
                    val = _extract_select_value(raw)
                    if val:
                        result[en_name] = category_map.get(val, "other")
                        result["_bt_category"] = val
                elif cn_name in ("制度状态", "法规状态"):
                    val = _extract_select_value(raw)
                    if val:
                        result[en_name] = status_map.get(val, "draft")
                elif cn_name in ("核心要点总结",):
                    val = _extract_rich_text(raw)
                    if val.strip():
                        result[en_name] = val.strip()
                else:
                    val = _extract_rich_text(raw)
                    if val.strip():
                        result[en_name] = val.strip()

            # 版本号
            if version_field:
                version_raw = bitable_fields.get(str(version_field))
                if version_raw:
                    version_text = _extract_rich_text(version_raw).strip()
                    if version_text:
                        result["_version_label"] = version_text
                        try:
                            if version_text in "ABCDEFGH":
                                result["version"] = ord(version_text) - ord("A") + 1
                            else:
                                result["version"] = int(version_text)
                        except (ValueError, TypeError):
                            result["version"] = 1

            # 附件（只存元数据）
            attachments = _extract_attachment_list(bitable_fields.get(attachment_field))
            if attachments:
                att = attachments[0]
                result["_attachment_file_token"] = str(att.get("file_token", ""))
                result["_attachment_original_name"] = str(att.get("name", ""))

            # 日期字段
            publish_date = _ms_timestamp_to_date(bitable_fields.get("颁布修订日期"))
            if publish_date:
                result["publish_date"] = publish_date
            impl_date = _ms_timestamp_to_date(bitable_fields.get("实施日期"))
            if impl_date:
                result["implementation_date"] = impl_date

            # 数据来源表
            result["_source"] = str(table_def["source"])

            return result

        async def generate_article_no(
            category: str, table_def: dict[str, object]
        ) -> str:
            prefix_map = table_def.get("article_no_prefix", default_prefix)
            prefix = prefix_map.get(category, "GEN")
            today = date.today().strftime("%Y%m%d")
            pattern = f"{prefix}-{today}-%"

            stmt = select(func.max(SafetyKnowledgeArticle.article_no)).where(
                SafetyKnowledgeArticle.article_no.ilike(pattern),
            )
            result = await self.session.execute(stmt)
            max_no = result.scalar()

            if max_no and str(max_no).startswith(f"{prefix}-{today}-"):
                try:
                    seq = int(str(max_no).rsplit("-", 1)[-1]) + 1
                except (ValueError, IndexError):
                    seq = 1
            else:
                seq = 1

            return f"{prefix}-{today}-{seq:03d}"

        # ── 1. 拉取 Bitable 全部记录 ──
        bitable_records: list[dict[str, object]] = []

        for td in tables:
            table_id = str(td["table_id"])
            table_name = str(td["name"])
            client = SafetyBitableClient(app_token="IkYTw6PJPiKZTCkfuNQc7IqEnzd", table_id=table_id)
            records = await client.list_all_records()
            logger.info("Bitable [%s]: %d 条", table_name, len(records))
            for rec in records:
                rid = rec.get("record_id", "")
                if rid:
                    bitable_records.append({
                        "record_id": str(rid),
                        "fields": rec.get("fields", {}),
                        "table_def": td,
                    })

        bitable_map: dict[str, dict[str, object]] = {
            str(r["record_id"]): r for r in bitable_records
        }
        total_bitable = len(bitable_map)

        # ── 2. 拉取平台全部记录 ──
        stmt = select(SafetyKnowledgeArticle).where(
            SafetyKnowledgeArticle.feishu_record_id.isnot(None),
            SafetyKnowledgeArticle.is_deleted == False,  # noqa: E712
        )
        result = await self.session.execute(stmt)
        platform_articles = result.scalars().all()
        platform_map: dict[str, SafetyKnowledgeArticle] = {}
        for a in platform_articles:
            if a.feishu_record_id:
                platform_map[a.feishu_record_id] = a

        bitable_ids = set(bitable_map.keys())
        platform_ids = set(platform_map.keys())

        to_create_ids = bitable_ids - platform_ids
        to_update_ids = bitable_ids & platform_ids
        to_delete_ids = platform_ids - bitable_ids

        logger.info(
            "Sync: create=%d update=%d delete=%d",
            len(to_create_ids), len(to_update_ids), len(to_delete_ids),
        )

        created_count = 0
        updated_count = 0
        deleted_count = 0

        # ── 3a. CREATE ──
        for rid in sorted(to_create_ids):
            rec = bitable_map[rid]
            td = rec["table_def"]
            fields = rec["fields"]
            mapped = map_fields(fields, td)

            title = str(mapped.get("title", ""))
            if not title:
                continue

            category = str(mapped.get("category", "other"))
            article_no = await generate_article_no(category, td)

            # article_no 冲突处理
            existing = await self.session.execute(
                select(SafetyKnowledgeArticle).where(
                    SafetyKnowledgeArticle.article_no == article_no,
                ),
            )
            conflict = existing.scalar_one_or_none()
            if conflict is not None:
                new_no = f"{conflict.article_no}-OLD"
                logger.warning("article_no conflict: %s → rename old to %s", article_no, new_no)
                await self.session.execute(
                    update(SafetyKnowledgeArticle)
                    .where(SafetyKnowledgeArticle.id == conflict.id)
                    .values(article_no=new_no, updated_at=datetime.now()),
                )
                await self.session.flush()

            # 组装 notes
            notes_parts: list[str] = []
            version_label = str(mapped.pop("_version_label", "") or "")
            if version_label:
                notes_parts.append(f"版本: {version_label}")
            if mapped.get("notes"):
                notes_parts.append(str(mapped["notes"]))

            bt_category = str(mapped.pop("_bt_category", "") or "")
            source = str(mapped.pop("_source", str(td.get("source", ""))))

            attachment_token = str(mapped.pop("_attachment_file_token", "") or "")
            attachment_name = str(mapped.pop("_attachment_original_name", "") or "")

            article = SafetyKnowledgeArticle(
                feishu_record_id=rid,
                article_no=article_no,
                title=title,
                category=category,
                tags=bt_category or None,
                status=str(mapped.get("status", "published")),
                source=source,
                author=str(mapped.get("author", "")) or None,
                content=str(mapped.get("content", "")) or None,
                publish_date=mapped.get("publish_date"),
                implementation_date=mapped.get("implementation_date"),
                notes="; ".join(notes_parts) if notes_parts else None,
                version=int(str(mapped.get("version", 1))),
                attachment_path=attachment_token or None,
                attachment_original_name=attachment_name or None,
            )
            self.session.add(article)
            await self.session.flush()
            created_count += 1

        # ── 3b. UPDATE ──
        for rid in sorted(to_update_ids):
            rec = bitable_map[rid]
            td = rec["table_def"]
            fields = rec["fields"]
            mapped = map_fields(fields, td)
            a = platform_map[rid]

            update_data: dict[str, object] = {}

            bt_category = str(mapped.pop("_bt_category", "") or "")
            if bt_category and bt_category != str(getattr(a, "tags", "") or ""):
                update_data["tags"] = bt_category

            source_val = str(mapped.pop("_source", ""))
            if source_val:
                mapped["source"] = source_val

            for en_name in ("title", "category", "status", "notes", "author", "content", "source"):
                new_val = mapped.get(en_name)
                old_val = getattr(a, en_name, None)
                if new_val is not None and str(new_val) != str(old_val or ""):
                    update_data[en_name] = new_val

            for date_field in ("publish_date", "implementation_date"):
                new_val = mapped.get(date_field)
                old_val = getattr(a, date_field, None)
                if new_val is not None and str(new_val) != str(old_val):
                    update_data[date_field] = new_val

            new_version = mapped.get("version")
            if new_version is not None and int(str(new_version)) != (getattr(a, "version", 1) or 1):
                update_data["version"] = new_version

            new_token = str(mapped.pop("_attachment_file_token", "") or "")
            if new_token and new_token != str(a.attachment_path or ""):
                update_data["attachment_path"] = new_token
            new_name = str(mapped.pop("_attachment_original_name", "") or "")
            if new_name and new_name != str(a.attachment_original_name or ""):
                update_data["attachment_original_name"] = new_name

            version_label = str(mapped.pop("_version_label", "") or "")
            if version_label:
                current_notes = str(getattr(a, "notes", "") or "")
                label_prefix = f"版本: {version_label}"
                if label_prefix not in current_notes:
                    new_notes = current_notes + ("; " + label_prefix if current_notes else label_prefix)
                    if new_notes != current_notes:
                        update_data["notes"] = new_notes

            if not update_data:
                continue

            await self.session.execute(
                update(SafetyKnowledgeArticle)
                .where(SafetyKnowledgeArticle.id == a.id)
                .values(**update_data, updated_at=datetime.now()),
            )
            updated_count += 1

        # ── 3c. DELETE ──
        for rid in sorted(to_delete_ids):
            a = platform_map[rid]
            if a.is_deleted:
                continue
            await self.session.execute(
                update(SafetyKnowledgeArticle)
                .where(SafetyKnowledgeArticle.id == a.id)
                .values(is_deleted=True, updated_at=datetime.now()),
            )
            deleted_count += 1

        # ── 4. 最终计数 ──
        count_stmt = select(func.count(SafetyKnowledgeArticle.id)).where(
            SafetyKnowledgeArticle.is_deleted == False,  # noqa: E712
            SafetyKnowledgeArticle.feishu_record_id.isnot(None),
        )
        count_result = await self.session.execute(count_stmt)
        total_platform = count_result.scalar() or 0

        return KnowledgeSyncResponse(
            created=created_count,
            updated=updated_count,
            deleted=deleted_count,
            total_bitable=total_bitable,
            total_platform=total_platform,
        )
